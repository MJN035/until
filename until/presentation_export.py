"""Until 발표 초안을 실제 PowerPoint(.pptx)로 내보낸다."""
from __future__ import annotations

from io import BytesIO
import re

_SLIDE = re.compile(r"^##\s*슬라이드\s*\d+\s*[:：-]?\s*(.*)$", re.I)
_DECISION = re.compile(r"\[\[DECISION:\s*(.*?)\]\]")


def parse_slide_markdown(text: str) -> list[tuple[str, list[str]]]:
    """Until의 슬라이드 Markdown을 제목/불릿 배열로 바꾼다."""
    slides: list[tuple[str, list[str]]] = []
    title = ""
    bullets: list[str] = []
    for raw in (text or "").splitlines():
        line = raw.strip()
        match = _SLIDE.match(line)
        if match:
            if title:
                slides.append((title, bullets))
            title, bullets = (match.group(1).strip() or "발표"), []
            continue
        if not title or not line or line.startswith("#"):
            continue
        line = re.sub(r"^[-*+]\s+|^\d+[.)]\s+", "", line)
        line = _DECISION.sub(lambda m: f"【직접 정할 것】 {m.group(1)}", line)
        line = re.sub(r"\*\*(.*?)\*\*|__(.*?)__", lambda m: m.group(1) or m.group(2), line)
        if line:
            bullets.append(line)
    if title:
        slides.append((title, bullets))
    return slides or [("발표 자료", ["초안 내용을 확인해 주세요."])]


def render_presentation_pptx(result) -> bytes:
    """Result의 완성본(없으면 초안)을 읽기 쉬운 16:9 PPTX bytes로 만든다.

    python-pptx는 발표 내보내기 표면에서만 지연 import한다. 따라서 mock 파이프라인과
    나머지 오프라인 기능은 이 선택 의존성 없이 계속 동작한다.
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError(
            "PPTX 내보내기 구성요소가 없습니다: pip install python-pptx"
        ) from exc

    draft = getattr(result, "final_draft", None) or result.draft
    slides = parse_slide_markdown(getattr(draft, "body", ""))
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    bg = RGBColor(12, 18, 32)
    title_color = RGBColor(246, 247, 251)
    body_color = RGBColor(220, 227, 240)
    accent = RGBColor(108, 140, 255)
    decision = RGBColor(86, 214, 179)

    for index, (title, bullets) in enumerate(slides, 1):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg

        bar = slide.shapes.add_shape(1, Inches(0.72), Inches(0.72),
                                     Inches(0.08), Inches(0.72))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.98), Inches(0.62),
                                              Inches(11.4), Inches(0.9))
        title_frame = title_box.text_frame
        title_frame.clear(); title_frame.word_wrap = True
        run = title_frame.paragraphs[0].add_run(); run.text = title[:90]
        run.font.name = "맑은 고딕"; run.font.size = Pt(30); run.font.bold = True
        run.font.color.rgb = title_color

        body_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.7),
                                             Inches(11.1), Inches(4.9))
        frame = body_box.text_frame
        frame.clear(); frame.word_wrap = True; frame.vertical_anchor = MSO_ANCHOR.TOP
        clean = [b for b in bullets if b][:5] or ["내용을 확인해 주세요."]
        total = sum(len(b) for b in clean)
        font_size = 22 if total <= 260 else 18
        for i, value in enumerate(clean):
            para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            para.text = value[:500]
            para.level = 0
            para.space_after = Pt(13)
            para.font.name = "맑은 고딕"; para.font.size = Pt(font_size)
            para.font.color.rgb = decision if value.startswith("【직접 정할 것】") else body_color

        num = slide.shapes.add_textbox(Inches(11.85), Inches(6.83),
                                        Inches(0.65), Inches(0.3))
        p = num.text_frame.paragraphs[0]; p.text = f"{index:02d}"
        p.font.name = "Aptos"; p.font.size = Pt(10); p.font.color.rgb = accent

    out = BytesIO()
    prs.save(out)
    return out.getvalue()
